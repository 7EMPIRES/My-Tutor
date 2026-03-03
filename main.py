import os
import sys
import base64
import io
import json
from typing import List, Optional, Union
from fastapi import FastAPI, Request, UploadFile, File, Form, HTTPException
from fastapi.responses import HTMLResponse, JSONResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
import google.generativeai as genai
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
from dotenv import load_dotenv

load_dotenv()

app = FastAPI()

# Helper for PyInstaller path resolution
def get_resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller """
    try:
        # PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

# Configure Gemini
api_key = os.getenv("GEMINI_API_KEY")
if api_key:
    genai.configure(api_key=api_key)

templates = Jinja2Templates(directory=get_resource_path("templates"))

# Mount static files
app.mount("/static", StaticFiles(directory=get_resource_path("static")), name="static")

@app.get("/api/health")
async def health():
    return {"status": "ok"}

@app.get("/sw.js")
async def serve_sw():
    return FileResponse(get_resource_path("static/sw.js"), media_type="application/javascript")

# Ensure static directory exists (for dev)
if not os.path.exists("static"):
    os.makedirs("static")

# Helper: Extract text and images from PDF
def extract_pdf_content(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = ""
    images = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text += page.get_text() + "\n"
        
        # Extract images
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            images.append(base64.b64encode(image_bytes).decode("utf-8"))
            
    return text, images

# Helper: Extract text and images from Word
def extract_docx_content(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    text = "\n".join([p.text for p in doc.paragraphs])
    images = []
    
    # Extract images from docx
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            image_bytes = rel.target_part.blob
            images.append(base64.b64encode(image_bytes).decode("utf-8"))
            
    return text, images

# Helper: Extract text and images from PPTX
def extract_pptx_content(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    images = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            if shape.shape_type == 13: # Picture
                image_bytes = shape.image.blob
                images.append(base64.b64encode(image_bytes).decode("utf-8"))
                
    return text, images

@app.get("/", response_class=HTMLResponse)
async def index(request: Request):
    return templates.TemplateResponse("index.html", {
        "request": request,
        "gemini_api_key": os.getenv("GEMINI_API_KEY")
    })

@app.post("/api/extract")
async def extract(
    file: Optional[UploadFile] = File(None)
):
    if not file:
        raise HTTPException(status_code=400, detail="No file uploaded")
    
    content = await file.read()
    filename = file.filename.lower()
    
    text = ""
    images = []
    
    try:
        if filename.endswith(".pdf"):
            text, images = extract_pdf_content(content)
        elif filename.endswith(".docx") or filename.endswith(".doc"):
            text, images = extract_docx_content(content)
        elif filename.endswith(".pptx") or filename.endswith(".ppt"):
            text, images = extract_pptx_content(content)
        elif filename.endswith((".jpg", ".jpeg", ".png", ".webp")):
            text = "Image file uploaded."
            images = [base64.b64encode(content).decode("utf-8")]
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Extraction failed: {str(e)}")

    return {
        "text": text,
        "images": images,
        "filename": file.filename
    }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=3000)
